
import logging
import hashlib
from abc import ABC, abstractmethod
from dataclasses import dataclass
from typing import Dict, List, Any, Optional, Union
from pathlib import Path
import mimetypes
# import magic

# PDF parsing
import PyPDF2
import pdfplumber

# Office parsing  
from docx import Document as DocxDocument
from pptx import Presentation

from src.utils.exceptions import DocumentProcessingException as DocumentParsingException

logger = logging.getLogger(__name__)

@dataclass
class DocumentBlock:
    """Unidad atómica de contenido estructurado"""
    content: str
    block_type: str  # "title", "paragraph", "list", "table", "metadata"  
    level: int       # Jerarquía (H1=1, H2=2, etc.)
    metadata: Dict[str, Any]

@dataclass
class ParsedDocument:
    """Resultado completo del parsing"""
    blocks: List[DocumentBlock]
    metadata: Dict[str, Any]
    source_info: Dict[str, Any]

class DocumentParser(ABC):
    """Contrato base para parsers especializados"""
    
    @abstractmethod
    def can_parse(self, file_path: str, mime_type: str) -> bool:
        """Verifica si puede procesar el formato"""
        pass
    
    @abstractmethod
    def parse(self, file_path: str) -> ParsedDocument:
        """Extrae contenido estructurado del documento"""  
        pass
    
    @abstractmethod
    def validate_document(self, file_path: str) -> Dict[str, Any]:
        """Valida integridad y seguridad del documento"""
        pass

class DocumentParserFactory:
    """Factory para crear parsers según formato de archivo"""
    
    _parsers: Dict[str, type] = {}
    
    @classmethod
    def register_parser(cls, extensions: List[str], parser_class: type):
        """Registra un parser para extensiones específicas"""
        for ext in extensions:
            cls._parsers[ext.lower().lstrip('.')] = parser_class
    
    @classmethod
    def get_parser(cls, file_path: str) -> DocumentParser:
        """Obtiene parser apropiado para el archivo"""
        extension = Path(file_path).suffix.lower().lstrip('.')
        
        if extension not in cls._parsers:
            raise DocumentParsingException(f"No parser available for extension: {extension}")
        
        parser_class = cls._parsers[extension]
        return parser_class()
    
    @classmethod
    def supported_formats(cls) -> List[str]:
        """Lista formatos soportados"""
        return list(cls._parsers.keys())

class PDFParser(DocumentParser):
    """Parser especializado para documentos PDF"""
    
    def can_parse(self, file_path: str, mime_type: str) -> bool:
        """Verifica si es PDF válido"""
        return mime_type == 'application/pdf' or file_path.lower().endswith('.pdf')
    
    def validate_document(self, file_path: str) -> Dict[str, Any]:
        """Valida integridad del PDF"""
        try:
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                return {
                    "valid": True,
                    "pages": len(reader.pages),
                    "encrypted": reader.is_encrypted,
                    "metadata": reader.metadata or {}
                }
        except Exception as e:
            return {"valid": False, "error": str(e)}
    
    def parse(self, file_path: str) -> ParsedDocument:
        """Extrae contenido estructurado del PDF"""
        try:
            blocks = []
            
            # Usar pdfplumber para mejor extracción de estructura
            with pdfplumber.open(file_path) as pdf:
                for page_num, page in enumerate(pdf.pages, 1):
                    
                    # Extraer texto con posición para detectar títulos
                    chars = page.chars
                    if not chars:
                        continue
                    
                    # Agrupar por líneas de texto
                    lines = self._group_chars_by_line(chars)
                    
                    for line_info in lines:
                        text = line_info['text'].strip()
                        if not text:
                            continue
                        
                        # Detectar tipo de bloque por características tipográficas
                        block_type, level = self._classify_pdf_text(line_info)
                        
                        blocks.append(DocumentBlock(
                            content=text,
                            block_type=block_type,
                            level=level,
                            metadata={
                                "page": page_num,
                                "font_size": line_info.get('font_size', 0),
                                "is_bold": line_info.get('is_bold', False),
                                "position": line_info.get('position', {})
                            }
                        ))
                    
                    # Extraer tablas
                    tables = page.extract_tables()
                    for table_idx, table in enumerate(tables or []):
                        if table:
                            table_text = self._format_table_as_text(table)
                            blocks.append(DocumentBlock(
                                content=table_text,
                                block_type="table",
                                level=0,
                                metadata={
                                    "page": page_num,
                                    "table_index": table_idx,
                                    "rows": len(table),
                                    "cols": len(table[0]) if table else 0
                                }
                            ))
            
            # Extraer metadatos del documento
            metadata = self._extract_pdf_metadata(file_path)
            
            return ParsedDocument(
                blocks=blocks,
                metadata=metadata,
                source_info={"file_path": file_path, "pages": len(pdf.pages)}
            )
            
        except Exception as e:
            logger.error(f"Error parsing PDF {file_path}: {e}")
            raise DocumentParsingException(f"Failed to parse PDF: {e}")
    
    def _group_chars_by_line(self, chars: List[Dict]) -> List[Dict]:
        """Agrupa caracteres por líneas de texto"""
        lines = []
        if not chars:
            return lines
        
        # Agrupar por posición Y similar (misma línea)
        current_line = {"chars": [], "y": chars[0]['y0']}
        
        for char in chars:
            # Si está en la misma línea (tolerancia de 2 puntos)
            if abs(char['y0'] - current_line['y']) <= 2:
                current_line['chars'].append(char)
            else:
                # Finalizar línea anterior
                if current_line['chars']:
                    lines.append(self._consolidate_line(current_line))
                
                # Comenzar nueva línea
                current_line = {"chars": [char], "y": char['y0']}
        
        # Agregar última línea
        if current_line['chars']:
            lines.append(self._consolidate_line(current_line))
        
        return lines
    
    def _consolidate_line(self, line_info: Dict) -> Dict:
        """Consolida caracteres de una línea en texto"""
        chars = line_info['chars']
        text = ''.join(char['text'] for char in chars)
        
        # Detectar propiedades tipográficas predominantes
        font_sizes = [char.get('size', 0) for char in chars]
        font_names = [char.get('fontname', '') for char in chars]
        
        avg_font_size = sum(font_sizes) / len(font_sizes) if font_sizes else 0
        is_bold = any('bold' in name.lower() for name in font_names)
        
        return {
            'text': text,
            'font_size': avg_font_size,
            'is_bold': is_bold,
            'position': {
                'x0': min(char['x0'] for char in chars),
                'y0': line_info['y'],
                'x1': max(char['x1'] for char in chars),
                'y1': max(char['y1'] for char in chars)
            }
        }
    
    def _classify_pdf_text(self, line_info: Dict) -> tuple:
        """Clasifica tipo de bloque por características tipográficas"""
        text = line_info['text']
        font_size = line_info.get('font_size', 0)
        is_bold = line_info.get('is_bold', False)
        
        # Detectar títulos por tamaño de fuente y estilo
        if font_size > 16 or (font_size > 14 and is_bold):
            return "title", 1
        elif font_size > 14 or (font_size > 12 and is_bold):
            return "title", 2
        elif font_size > 12 or is_bold:
            return "title", 3
        
        # Detectar listas
        if text.lstrip().startswith(('•', '-', '*', '1.', '2.', '3.')):
            return "list", 0
        
        return "paragraph", 0
    
    def _format_table_as_text(self, table: List[List[str]]) -> str:
        """Convierte tabla a texto estructurado"""
        if not table:
            return ""
        
        # Crear texto de tabla con separadores
        lines = []
        for row in table:
            clean_row = [str(cell or '').strip() for cell in row]
            lines.append(' | '.join(clean_row))
        
        return '\n'.join(lines)
    
    def _extract_pdf_metadata(self, file_path: str) -> Dict[str, Any]:
        """Extrae metadatos del PDF"""
        try:
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                metadata = reader.metadata or {}
                
                return {
                    "title": metadata.get('/Title', ''),
                    "author": metadata.get('/Author', ''),
                    "subject": metadata.get('/Subject', ''),
                    "creator": metadata.get('/Creator', ''),
                    "producer": metadata.get('/Producer', ''),
                    "creation_date": str(metadata.get('/CreationDate', '')),
                    "modification_date": str(metadata.get('/ModDate', '')),
                    "pages": len(reader.pages)
                }
        except Exception:
            return {}

class DOCXParser(DocumentParser):
    """Parser especializado para documentos DOCX"""
    
    def can_parse(self, file_path: str, mime_type: str) -> bool:
        """Verifica si es DOCX válido"""
        return (mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document' 
                or file_path.lower().endswith('.docx'))
    
    def validate_document(self, file_path: str) -> Dict[str, Any]:
        """Valida integridad del DOCX"""
        try:
            doc = DocxDocument(file_path)
            return {
                "valid": True,
                "paragraphs": len(doc.paragraphs),
                "tables": len(doc.tables),
                "sections": len(doc.sections)
            }
        except Exception as e:
            return {"valid": False, "error": str(e)}
    
    def parse(self, file_path: str) -> ParsedDocument:
        """Extrae contenido estructurado del DOCX"""
        try:
            doc = DocxDocument(file_path)
            blocks = []
            
            # Procesar párrafos con estilos
            for para in doc.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                
                # Detectar tipo y nivel por estilo
                block_type, level = self._classify_docx_paragraph(para)
                
                blocks.append(DocumentBlock(
                    content=text,
                    block_type=block_type,
                    level=level,
                    metadata={
                        "style": para.style.name,
                        "is_bold": any(run.bold for run in para.runs),
                        "is_italic": any(run.italic for run in para.runs),
                        "font_size": self._get_paragraph_font_size(para)
                    }
                ))
            
            # Procesar tablas
            for table_idx, table in enumerate(doc.tables):
                table_text = self._extract_docx_table(table)
                blocks.append(DocumentBlock(
                    content=table_text,
                    block_type="table", 
                    level=0,
                    metadata={
                        "table_index": table_idx,
                        "rows": len(table.rows),
                        "cols": len(table.columns)
                    }
                ))
            
            metadata = self._extract_docx_metadata(doc)
            
            return ParsedDocument(
                blocks=blocks,
                metadata=metadata,
                source_info={"file_path": file_path, "paragraphs": len(doc.paragraphs)}
            )
            
        except Exception as e:
            logger.error(f"Error parsing DOCX {file_path}: {e}")
            raise DocumentParsingException(f"Failed to parse DOCX: {e}")
    
    def _classify_docx_paragraph(self, para) -> tuple:
        """Clasifica párrafo por estilo Word"""
        style_name = para.style.name.lower()
        
        # Mapear estilos de Word a tipos
        if 'heading 1' in style_name or 'título 1' in style_name:
            return "title", 1
        elif 'heading 2' in style_name or 'título 2' in style_name:
            return "title", 2
        elif 'heading 3' in style_name or 'título 3' in style_name:
            return "title", 3
        elif 'heading' in style_name or 'título' in style_name:
            return "title", 4
        elif 'list' in style_name or 'bullet' in style_name:
            return "list", 0
        
        return "paragraph", 0
    
    def _get_paragraph_font_size(self, para) -> Optional[float]:
        """Obtiene tamaño de fuente predominante del párrafo"""
        sizes = []
        for run in para.runs:
            if run.font.size:
                sizes.append(run.font.size.pt)
        
        return sum(sizes) / len(sizes) if sizes else None
    
    def _extract_docx_table(self, table) -> str:
        """Extrae texto de tabla DOCX"""
        rows = []
        for row in table.rows:
            cells = []
            for cell in row.cells:
                cell_text = ' '.join(para.text.strip() for para in cell.paragraphs).strip()
                cells.append(cell_text)
            rows.append(' | '.join(cells))
        
        return '\n'.join(rows)
    
    def _extract_docx_metadata(self, doc) -> Dict[str, Any]:
        """Extrae metadatos del DOCX"""
        props = doc.core_properties
        return {
            "title": props.title or "",
            "author": props.author or "",
            "subject": props.subject or "",
            "keywords": props.keywords or "",
            "comments": props.comments or "",
            "created": str(props.created) if props.created else "",
            "modified": str(props.modified) if props.modified else ""
        }

class PPTXParser(DocumentParser):
    """Parser especializado para presentaciones PPTX"""
    
    def can_parse(self, file_path: str, mime_type: str) -> bool:
        """Verifica si es PPTX válido"""
        return (mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
                or file_path.lower().endswith('.pptx'))
    
    def validate_document(self, file_path: str) -> Dict[str, Any]:
        """Valida integridad del PPTX"""
        try:
            prs = Presentation(file_path)
            return {
                "valid": True,
                "slides": len(prs.slides),
                "slide_layouts": len(prs.slide_layouts)
            }
        except Exception as e:
            return {"valid": False, "error": str(e)}
    
    def parse(self, file_path: str) -> ParsedDocument:
        """Extrae contenido estructurado del PPTX"""
        try:
            prs = Presentation(file_path)
            blocks = []
            
            # Procesar cada slide
            for slide_num, slide in enumerate(prs.slides, 1):
                
                # Extraer título del slide
                if slide.shapes.title:
                    title_text = slide.shapes.title.text.strip()
                    if title_text:
                        blocks.append(DocumentBlock(
                            content=title_text,
                            block_type="title",
                            level=1,
                            metadata={
                                "slide": slide_num,
                                "element": "slide_title"
                            }
                        ))
                
                # Extraer contenido de placeholders
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        if shape != slide.shapes.title:  # No duplicar título
                            blocks.append(DocumentBlock(
                                content=shape.text.strip(),
                                block_type="paragraph",
                                level=2,
                                metadata={
                                    "slide": slide_num,
                                    "shape_type": str(shape.shape_type),
                                    "element": "content"
                                }
                            ))
                
                # Extraer tablas
                for shape in slide.shapes:
                    if shape.has_table:
                        table_text = self._extract_pptx_table(shape.table)
                        blocks.append(DocumentBlock(
                            content=table_text,
                            block_type="table",
                            level=0,
                            metadata={
                                "slide": slide_num,
                                "rows": len(shape.table.rows),
                                "cols": len(shape.table.columns)
                            }
                        ))
            
            metadata = self._extract_pptx_metadata(prs)
            
            return ParsedDocument(
                blocks=blocks,
                metadata=metadata,
                source_info={"file_path": file_path, "slides": len(prs.slides)}
            )
            
        except Exception as e:
            logger.error(f"Error parsing PPTX {file_path}: {e}")
            raise DocumentParsingException(f"Failed to parse PPTX: {e}")
    
    def _extract_pptx_table(self, table) -> str:
        """Extrae texto de tabla PPTX"""
        rows = []
        for row in table.rows:
            cells = []
            for cell in row.cells:
                cell_text = cell.text.strip()
                cells.append(cell_text)
            rows.append(' | '.join(cells))
        
        return '\n'.join(rows)
    
    def _extract_pptx_metadata(self, prs) -> Dict[str, Any]:
        """Extrae metadatos del PPTX"""
        props = prs.core_properties
        return {
            "title": props.title or "",
            "author": props.author or "",
            "subject": props.subject or "",
            "keywords": props.keywords or "",
            "comments": props.comments or "",
            "created": str(props.created) if props.created else "",
            "modified": str(props.modified) if props.modified else ""
        }

# Registro automático de parsers
DocumentParserFactory.register_parser(['pdf'], PDFParser)
DocumentParserFactory.register_parser(['docx'], DOCXParser)
DocumentParserFactory.register_parser(['pptx'], PPTXParser)